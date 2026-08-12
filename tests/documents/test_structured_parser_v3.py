from __future__ import annotations

import asyncio
import io

import fitz

from backend.app.services.document_parsing.models import (
    DocumentBlockType,
    DocumentParsingProfile,
    OCRResult,
)
from backend.app.services.document_parsing.parser import DocumentParser


def _pdf_with_layout() -> bytes:
    document = fitz.open()
    page = document.new_page(width=600, height=800)
    page.insert_text((60, 80), "RAG Architecture", fontsize=22, fontname="helv")
    page.insert_textbox(
        fitz.Rect(60, 120, 540, 220),
        "Retrieval combines structure and semantic evidence. This paragraph is intentionally long enough to remain native text.",
        fontsize=11,
    )
    page.draw_rect(fitz.Rect(60, 260, 300, 330), color=(0, 0, 0), fill=None)
    page.draw_line((60, 285), (300, 285), color=(0, 0, 0))
    page.draw_line((180, 260), (180, 330), color=(0, 0, 0))
    page.insert_text((75, 278), "metric", fontsize=10)
    page.insert_text((195, 278), "score", fontsize=10)
    page.insert_text((75, 310), "Recall", fontsize=10)
    page.insert_text((195, 310), "0.91", fontsize=10)
    content = document.tobytes()
    document.close()
    return content


def test_pdf_structured_profile_emits_heading_paragraph_and_version() -> None:
    result = asyncio.run(
        DocumentParser().parse_document(
            content=_pdf_with_layout(),
            filename="lesson.pdf",
            mime_type="application/pdf",
            profile=DocumentParsingProfile.STRUCTURED_V3,
        )
    )

    assert result.parser_version == "document-parser-v4"
    assert any(block.block_type is DocumentBlockType.HEADING for block in result.blocks)
    assert any(block.block_type is DocumentBlockType.PARAGRAPH for block in result.blocks)
    assert all(block.bbox is not None for block in result.blocks)
    assert all(block.reading_order is not None for block in result.blocks)
    assert all(
        block.source_char_start is not None and block.source_char_end is not None
        for block in result.blocks
    )


def test_pdf_structured_profile_emits_valid_table_without_duplicate_text() -> None:
    result = asyncio.run(
        DocumentParser().parse_document(
            content=_pdf_with_layout(),
            filename="table.pdf",
            mime_type="application/pdf",
            profile=DocumentParsingProfile.STRUCTURED_V3,
        )
    )

    tables = [block for block in result.blocks if block.block_type is DocumentBlockType.TABLE]
    assert len(tables) <= 1
    if tables:
        assert "metric" in tables[0].text
        assert "score" in tables[0].text
        assert tables[0].table_header_rows == 1
        assert sum("metric" in block.text for block in result.blocks) == 1


def test_pdf_structured_table_detector_degrades_to_paragraph(monkeypatch) -> None:
    monkeypatch.setattr("fitz.Page.find_tables", lambda *args, **kwargs: (_ for _ in ()).throw(RuntimeError("unavailable")))
    result = asyncio.run(
        DocumentParser().parse_document(
            content=_pdf_with_layout(),
            filename="fallback.pdf",
            mime_type="application/pdf",
            profile=DocumentParsingProfile.STRUCTURED_V3,
        )
    )

    assert result.status.value == "success"
    assert all(block.block_type is not DocumentBlockType.TABLE for block in result.blocks)
    assert any("metric" in block.text for block in result.blocks)


def test_spatial_table_detector_is_deterministic_fallback() -> None:
    from backend.app.services.document_parsing.pdf_parser import TableDetectionMethod, TableDetector

    def cell(text: str, x: float, y: float) -> dict:
        return {"bbox": (x, y, x + 40, y + 12), "lines": [{"bbox": (x, y, x + 40, y + 12), "spans": [{"text": text}]}]}

    class Page:
        def find_tables(self, **kwargs):
            raise RuntimeError("detector unavailable")

    tables = TableDetector().detect(
        Page(),
        text_blocks=[cell("header-a", 10, 10), cell("header-b", 60, 10), cell("value-a", 10, 30), cell("value-b", 60, 30)],
    )

    assert len(tables) == 1
    assert tables[0].method is TableDetectionMethod.SPATIAL_HEURISTIC
    assert tables[0].rows == (("header-a", "header-b"), ("value-a", "value-b"))


def test_pptx_structured_profile_uses_row_banding_and_preserves_lists() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(4), Inches(0.4))
    title.text = "Lesson title"
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3), Inches(0.4))
    left.text = "Left column"
    right = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(3), Inches(0.4))
    right.text = "Right column"
    body = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(0.8))
    body.text_frame.text = "First item"
    paragraph = body.text_frame.add_paragraph()
    paragraph.text = "Second item"
    paragraph.level = 1
    buffer = io.BytesIO()
    presentation.save(buffer)

    result = asyncio.run(
        DocumentParser().parse_document(
            content=buffer.getvalue(),
            filename="lesson.pptx",
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            profile=DocumentParsingProfile.STRUCTURED_V3,
        )
    )

    assert result.parser_version == "document-parser-v4"
    assert result.blocks[0].block_type is DocumentBlockType.SLIDE_TITLE
    texts = [block.text for block in result.blocks]
    assert texts.index("Left column") < texts.index("Right column")
    assert any(block.block_type is DocumentBlockType.LIST_ITEM for block in result.blocks)
    assert all(block.bbox is not None for block in result.blocks)


def test_legacy_profile_remains_page_or_slide_level() -> None:
    result = asyncio.run(
        DocumentParser().parse_document(
            content=_pdf_with_layout(),
            filename="legacy.pdf",
            mime_type="application/pdf",
        )
    )

    assert result.parser_version == "document-parser-v3"
    assert len(result.blocks) == 1
    assert result.blocks[0].block_type is DocumentBlockType.UNKNOWN
