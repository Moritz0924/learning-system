from __future__ import annotations

import io

import pytest
from PIL import Image
from pypdf import PdfWriter
from pptx import Presentation
from sqlalchemy import select

from backend.app.models import OutboxEvent
from backend.app.services.object_storage import ObjectStorageUnavailable
from tests.conftest import register_user_with_goal


def _pdf_bytes() -> bytes:
    output = io.BytesIO()
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    writer.write(output)
    return output.getvalue()


def _pptx_bytes() -> bytes:
    output = io.BytesIO()
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(output)
    return output.getvalue()


def _png_bytes() -> bytes:
    output = io.BytesIO()
    Image.new("RGB", (2, 2), "white").save(output, format="PNG")
    return output.getvalue()


@pytest.mark.parametrize(
    ("filename", "mime_type", "content", "expected_mime"),
    [
        ("guide.pdf", "application/pdf", _pdf_bytes(), "application/pdf"),
        (
            "slides.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            _pptx_bytes(),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
        ("diagram.png", "image/png", _png_bytes(), "image/png"),
        ("notes.md", "text/markdown", b"# Notes\nBounded upload.", "text/markdown"),
        ("notes.txt", "text/plain", b"Bounded plain text upload.", "text/plain"),
    ],
    ids=["pdf", "pptx", "png", "markdown", "text"],
)
def test_multipart_upload_accepts_real_supported_files(
    client, session_factory, monkeypatch, filename, mime_type, content, expected_mime
):
    monkeypatch.setenv("DOCUMENT_PROCESSING_MODE", "defer")
    account = register_user_with_goal(
        client,
        session_factory,
        email=f"{filename.replace('.', '-')}@example.com",
    )

    response = client.post(
        "/api/documents",
        headers=account["headers"],
        data={"goal_id": account["goal_id"]},
        files={"file": (filename, content, mime_type)},
    )

    assert response.status_code == 201, response.text
    body = response.json()
    assert body["filename"] == filename
    assert body["mime_type"] == expected_mime
    assert body["parse_status"] == "pending"
    assert "object_key" not in body
    assert "sha256" not in body


@pytest.mark.parametrize(
    ("filename", "mime_type", "content", "expected_status"),
    [
        ("../escape.txt", "text/plain", b"unsafe filename", 400),
        ("payload.exe", "application/octet-stream", b"MZ", 415),
        ("spoofed.pdf", "application/pdf", _png_bytes(), 415),
        ("spoofed.txt", "text/plain", _pdf_bytes(), 415),
        ("broken.pdf", "application/pdf", b"%PDF-this-is-not-a-pdf", 422),
        ("empty.txt", "text/plain", b"", 422),
    ],
    ids=[
        "unsafe-name",
        "unsupported",
        "mime-spoof",
        "extension-spoof",
        "corrupt",
        "empty",
    ],
)
def test_multipart_upload_rejects_invalid_or_spoofed_files(
    client, session_factory, monkeypatch, filename, mime_type, content, expected_status
):
    monkeypatch.setenv("DOCUMENT_PROCESSING_MODE", "defer")
    account = register_user_with_goal(
        client,
        session_factory,
        email=f"invalid-{expected_status}-{filename.replace('.', '-').replace('/', '-')}@example.com",
    )

    response = client.post(
        "/api/documents",
        headers=account["headers"],
        data={"goal_id": account["goal_id"]},
        files={"file": (filename, content, mime_type)},
    )

    assert response.status_code == expected_status, response.text


def test_multipart_upload_rejects_client_controlled_fields(
    client, session_factory, monkeypatch
):
    monkeypatch.setenv("DOCUMENT_PROCESSING_MODE", "defer")
    account = register_user_with_goal(
        client, session_factory, email="multipart-fields@example.com"
    )

    response = client.post(
        "/api/documents",
        headers=account["headers"],
        files={"file": ("notes.txt", b"safe text", "text/plain")},
        data={
            "goal_id": account["goal_id"],
            "user_id": "another-user",
            "trusted_level": "99",
        },
    )

    assert response.status_code == 422
    assert client.get("/api/documents", headers=account["headers"]).json() == {"documents": []}


def test_multipart_upload_enforces_decoded_file_limit(
    client, session_factory, monkeypatch
):
    monkeypatch.setenv("DOCUMENT_PROCESSING_MODE", "defer")
    monkeypatch.setenv("DOCUMENT_MAX_UPLOAD_BYTES", "4")
    account = register_user_with_goal(
        client, session_factory, email="multipart-limit@example.com"
    )

    response = client.post(
        "/api/documents",
        headers=account["headers"],
        data={"goal_id": account["goal_id"]},
        files={"file": ("notes.txt", b"12345", "text/plain")},
    )

    assert response.status_code == 413
    assert client.get("/api/documents", headers=account["headers"]).json() == {"documents": []}


def test_multipart_upload_returns_503_without_document_when_object_storage_fails(
    client, session_factory, monkeypatch
):
    class FailingStorage:
        def put_bytes(self, object_key, content, *, content_type):
            raise ObjectStorageUnavailable("storage unavailable")

        def delete_bytes(self, object_key):
            return None

    monkeypatch.setenv("DOCUMENT_PROCESSING_MODE", "defer")
    monkeypatch.setattr(
        "backend.app.application.document_service.build_document_object_storage",
        lambda: FailingStorage(),
    )
    account = register_user_with_goal(
        client, session_factory, email="multipart-storage@example.com"
    )

    response = client.post(
        "/api/documents",
        headers=account["headers"],
        data={"goal_id": account["goal_id"]},
        files={"file": ("notes.txt", b"safe content", "text/plain")},
    )

    assert response.status_code == 503
    assert client.get("/api/documents", headers=account["headers"]).json() == {"documents": []}


def test_inline_multipart_upload_returns_503_when_v3_embedding_is_unavailable(
    client, session_factory, monkeypatch
):
    monkeypatch.setenv("DOCUMENT_PROCESSING_MODE", "inline")
    monkeypatch.setenv("FEATURE_HYBRID_CHUNKING_V3", "true")
    monkeypatch.setenv("EMBEDDING_BACKEND", "openai")
    monkeypatch.delenv("EMBEDDING_API_KEY", raising=False)
    monkeypatch.delenv("LLM_API_KEY", raising=False)
    account = register_user_with_goal(
        client,
        session_factory,
        email="multipart-v3-embedding@example.com",
    )

    response = client.post(
        "/api/documents",
        headers=account["headers"],
        data={"goal_id": account["goal_id"]},
        files={"file": ("notes.md", b"# Notes\nNeeds embedding.", "text/markdown")},
    )

    assert response.status_code == 503
    assert "EMBEDDING_API_KEY" in response.json()["detail"]
    assert client.get("/api/documents", headers=account["headers"]).json() == {"documents": []}


def test_multipart_upload_keeps_pending_record_when_initial_celery_dispatch_fails(
    client, session_factory, monkeypatch
):
    import backend.app.worker as worker

    def fail_dispatch(*args, **kwargs):
        raise RuntimeError("broker secret detail")

    monkeypatch.setenv("DOCUMENT_PROCESSING_MODE", "celery")
    monkeypatch.setenv("DOCUMENT_PROCESSING_RETRY_DELAY_SECONDS", "0")
    monkeypatch.setattr(worker.process_document_upload_task, "delay", fail_dispatch)
    account = register_user_with_goal(
        client, session_factory, email="multipart-dispatch@example.com"
    )

    response = client.post(
        "/api/documents",
        headers=account["headers"],
        data={"goal_id": account["goal_id"]},
        files={"file": ("queued.txt", b"durable queue content", "text/plain")},
    )

    assert response.status_code == 201
    assert response.json()["parse_status"] == "pending"
    assert "broker secret detail" not in str(response.json())
    with session_factory() as session:
        event = session.scalar(
            select(OutboxEvent).where(
                OutboxEvent.payload_json["document_id"].as_string()
                == response.json()["id"]
            )
        )
        assert event is not None
        assert event.status == "pending"


def test_multipart_upload_isolated_by_authenticated_owner(
    client, session_factory, monkeypatch
):
    monkeypatch.setenv("DOCUMENT_PROCESSING_MODE", "defer")
    owner = register_user_with_goal(
        client, session_factory, email="multipart-owner@example.com"
    )
    other = register_user_with_goal(
        client, session_factory, email="multipart-other@example.com"
    )
    uploaded = client.post(
        "/api/documents",
        headers=owner["headers"],
        data={"goal_id": owner["goal_id"]},
        files={"file": ("private.txt", b"owner only", "text/plain")},
    )

    document_id = uploaded.json()["id"]
    assert client.get(f"/api/documents/{document_id}", headers=other["headers"]).status_code == 404
    assert client.get("/api/documents", headers=other["headers"]).json() == {"documents": []}


def test_legacy_markdown_upload_is_marked_deprecated_in_openapi(client):
    operation = client.get("/openapi.json").json()["paths"]["/api/documents/upload"]["post"]

    assert operation["deprecated"] is True
