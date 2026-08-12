from backend.app.services.document_parsing.models import (
    DocumentBlock,
    DocumentFileType,
    ProcessingMode,
    SourceElementType,
)
from backend.app.services.ocr import build_ocr_result_from_tesseract_data
from backend.app.services.document_parsing.fallback_policy import VisionFallbackPolicy
from backend.app.services.document_parsing.models import OCRResult
from backend.app.services.document_parsing.models import VisionContext, VisionEnrichmentStatus, VisionResult
from backend.app.services.document_parsing.parser import DocumentParser
from backend.app.services.vision_understanding import VisionClient

import httpx


def test_document_block_serializes_stable_parser_metadata():
    block = DocumentBlock(
        filename="lesson.pdf",
        file_type=DocumentFileType.PDF,
        page_number=3,
        block_index=2,
        text="OCR text",
        processing_mode=ProcessingMode.PDF_OCR,
        source_element=SourceElementType.PDF_RENDER,
        ocr_confidence=0.76,
    )

    assert block.model_dump(mode="json") == {
        "filename": "lesson.pdf",
        "file_type": "pdf",
        "page_number": 3,
        "block_index": 2,
        "text": "OCR text",
        "processing_mode": "pdf_ocr",
        "source_element": "pdf_render",
        "ocr_confidence": 0.76,
        "vision_confidence": None,
        "vision_enriched": False,
        "vision_enrichment_status": "not_needed",
        "source_element_index": None,
        "warnings": [],
    }


def test_ocr_result_ignores_invalid_words_and_weights_confidence_by_characters():
    result = build_ocr_result_from_tesseract_data(
        {
            "text": ["short", "longerword", "", "ignored"],
            "conf": ["50", "90", "80", "-1"],
            "left": [1, 2, 3, 4],
            "top": [1, 2, 3, 4],
            "width": [1, 2, 3, 4],
            "height": [1, 2, 3, 4],
        }
    )

    assert result.text == "short longerword"
    assert result.word_count == 2
    assert result.confidence == (0.5 * 5 + 0.9 * 10) / 15
    assert [word.text for word in result.words] == ["short", "longerword"]


def test_auto_vision_fallback_triggers_for_low_confidence_ocr(monkeypatch):
    monkeypatch.setenv("OCR_VISION_FALLBACK", "auto")
    monkeypatch.setenv("OCR_MIN_CONFIDENCE", "0.65")
    policy = VisionFallbackPolicy()

    assert policy.should_enrich(
        ocr_result=OCRResult(text="unclear", confidence=0.4, word_count=1, text_char_count=7),
        file_type=DocumentFileType.IMAGE,
        page_number=1,
    )


def test_document_parser_routes_an_image_to_structured_ocr_block():
    image_content = __import__("base64").b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR4nGP4////fwAJ+wP9KobjigAAAABJRU5ErkJggg=="
    )

    class FakeOCR:
        async def recognize_bytes(self, content: bytes, *, filename: str) -> OCRResult:
            assert filename == "slide.png"
            assert content == image_content
            return OCRResult(text="scanned lesson", confidence=0.8, word_count=2, text_char_count=14)

    result = __import__("asyncio").run(
        DocumentParser(ocr_service=FakeOCR()).parse_document(
            content=image_content, filename="slide.png", mime_type="image/png"
        )
    )

    assert result.status.value == "success"
    assert result.content_sha256
    assert result.blocks[0].processing_mode.value == "image_ocr"
    assert result.blocks[0].text == "scanned lesson"


def test_document_parser_extracts_pdf_text_layer_without_ocr():
    import fitz

    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "native PDF lesson text with sufficient content to avoid the OCR fallback path")
    content = document.tobytes()
    document.close()

    result = __import__("asyncio").run(
        DocumentParser().parse_document(content=content, filename="lesson.pdf", mime_type="application/pdf")
    )

    assert result.status.value == "success"
    assert [(block.processing_mode.value, block.page_number) for block in result.blocks] == [("pdf_text", 1)]
    assert "native PDF lesson text" in result.blocks[0].text


def test_document_parser_extracts_ppt_text_shapes_without_ocr():
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, 1_000_000, 1_000_000).text_frame.text = "native slide text"
    buffer = __import__("io").BytesIO()
    presentation.save(buffer)

    result = __import__("asyncio").run(
        DocumentParser().parse_document(
            content=buffer.getvalue(),
            filename="lesson.pptx",
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    )

    assert result.status.value == "success"
    assert [(block.processing_mode.value, block.page_number) for block in result.blocks] == [("ppt_native_text", 1)]
    assert result.blocks[0].text == "native slide text"


def test_document_parser_ocr_extracts_embedded_pptx_image():
    from PIL import Image
    from pptx import Presentation

    image = __import__("io").BytesIO()
    Image.new("RGB", (20, 20), "white").save(image, format="PNG")
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(__import__("io").BytesIO(image.getvalue()), 0, 0, 1_000_000, 1_000_000)
    content = __import__("io").BytesIO()
    presentation.save(content)

    class FakeOCR:
        async def recognize_bytes(self, image_content: bytes, *, filename: str) -> OCRResult:
            return OCRResult(text="scanned slide", confidence=0.9, word_count=2, text_char_count=13)

    result = __import__("asyncio").run(
        DocumentParser(ocr_service=FakeOCR()).parse_document(
            content=content.getvalue(), filename="scanned.pptx",
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    )

    assert [(block.processing_mode.value, block.source_element_index) for block in result.blocks] == [("ppt_ocr", 1)]


def test_document_parser_ocr_extracts_scanned_pdf_page():
    import fitz
    from PIL import Image

    image = __import__("io").BytesIO()
    Image.new("RGB", (50, 50), "white").save(image, format="PNG")
    document = fitz.open()
    page = document.new_page()
    page.insert_image(page.rect, stream=image.getvalue())
    content = document.tobytes()
    document.close()

    class FakeOCR:
        async def recognize_bytes(self, image_content: bytes, *, filename: str) -> OCRResult:
            return OCRResult(text="scanned PDF", confidence=0.9, word_count=2, text_char_count=11)

    result = __import__("asyncio").run(
        DocumentParser(ocr_service=FakeOCR()).parse_document(content=content, filename="scan.pdf", mime_type="application/pdf")
    )

    assert [(block.processing_mode.value, block.page_number) for block in result.blocks] == [("pdf_ocr", 1)]


def test_image_parser_persists_non_duplicate_vision_supplemental_text(monkeypatch):
    monkeypatch.setenv("OCR_VISION_FALLBACK", "always")
    image_content = __import__("base64").b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR4nGP4////fwAJ+wP9KobjigAAAABJRU5ErkJggg=="
    )

    class FakeOCR:
        async def recognize_bytes(self, content: bytes, *, filename: str) -> OCRResult:
            return OCRResult(text="OCR heading", confidence=0.9, word_count=2, text_char_count=11)

    class FakeVision:
        async def analyze_image(self, *args, **kwargs) -> VisionResult:
            return VisionResult(
                supplemental_text="OCR heading\nChart label",
                confidence=0.8,
                status=VisionEnrichmentStatus.SUCCESS,
            )

    result = __import__("asyncio").run(
        DocumentParser(ocr_service=FakeOCR(), vision_client=FakeVision()).parse_document(
            content=image_content, filename="diagram.png", mime_type="image/png"
        )
    )

    assert result.blocks[0].text == "OCR heading\nChart label"
    assert result.blocks[0].vision_enriched is True


def test_zhipu_vision_uses_independent_credentials_and_parses_fenced_json(monkeypatch):
    seen = {}
    monkeypatch.setenv("LLM_BASE_URL", "https://api.deepseek.com")
    monkeypatch.setenv("LLM_API_KEY", "deepseek-secret")
    monkeypatch.setenv("VISION_BASE_URL", "https://open.bigmodel.cn/api/paas/v4")
    monkeypatch.setenv("VISION_API_KEY", "zhipu-secret")
    monkeypatch.setenv("VISION_MODEL", "glm-4.5v")

    def handler(request: httpx.Request) -> httpx.Response:
        seen["url"] = str(request.url)
        seen["authorization"] = request.headers.get("authorization")
        seen["json"] = __import__("json").loads(request.content.decode())
        return httpx.Response(
            200,
            json={
                "choices": [
                    {
                        "message": {
                            "content": "```json\n{\"supplemental_text\":\"Chart label\",\"confidence\":0.9,\"complex_visual\":true}\n```"
                        }
                    }
                ]
            },
        )

    async def run():
        async with httpx.AsyncClient(transport=httpx.MockTransport(handler)) as http_client:
            return await VisionClient(http_client=http_client).analyze_image(
                b"image-bytes",
                mime_type="image/png",
                context=VisionContext(
                    filename="chart.png",
                    file_type=DocumentFileType.IMAGE,
                    page_number=1,
                    source_element=SourceElementType.IMAGE_FILE,
                ),
            )

    result = __import__("asyncio").run(run())

    assert result.status is VisionEnrichmentStatus.SUCCESS
    assert result.supplemental_text == "Chart label"
    assert result.confidence == 0.9
    assert result.complex_visual is True
    assert seen["url"] == "https://open.bigmodel.cn/api/paas/v4/chat/completions"
    assert seen["authorization"] == "Bearer zhipu-secret"
    assert seen["json"]["model"] == "glm-4.5v"
    assert seen["json"]["thinking"] == {"type": "enabled"}
    image_url = seen["json"]["messages"][1]["content"][1]["image_url"]["url"]
    assert image_url == __import__("base64").b64encode(b"image-bytes").decode("ascii")


def test_vision_client_does_not_fall_back_to_deepseek_credentials(monkeypatch):
    monkeypatch.setenv("LLM_BASE_URL", "https://api.deepseek.com")
    monkeypatch.setenv("LLM_API_KEY", "deepseek-secret")
    monkeypatch.delenv("VISION_API_KEY", raising=False)

    def handler(request: httpx.Request) -> httpx.Response:
        raise AssertionError("vision must not send an image to DeepSeek")

    async def run():
        async with httpx.AsyncClient(transport=httpx.MockTransport(handler)) as http_client:
            return await VisionClient(http_client=http_client).analyze_image(
                b"image-bytes",
                mime_type="image/png",
                context=VisionContext(
                    filename="chart.png",
                    file_type=DocumentFileType.IMAGE,
                    page_number=1,
                    source_element=SourceElementType.IMAGE_FILE,
                ),
            )

    result = __import__("asyncio").run(run())

    assert result.status is VisionEnrichmentStatus.UNAVAILABLE
    assert result.error_code == "vision_unavailable"


def test_vision_client_extracts_final_json_after_thinking_text(monkeypatch):
    monkeypatch.setenv("VISION_BASE_URL", "https://open.bigmodel.cn/api/paas/v4")
    monkeypatch.setenv("VISION_API_KEY", "zhipu-secret")
    monkeypatch.setenv("VISION_MODEL", "glm-4.5v")

    def handler(request: httpx.Request) -> httpx.Response:
        return httpx.Response(
            200,
            json={
                "choices": [
                    {
                        "message": {
                            "content": (
                                '<think>{"draft":"ignore"}</think>\n'
                                '{"supplemental_text":"Final chart label","confidence":0.8,"complex_visual":false}'
                            )
                        }
                    }
                ]
            },
        )

    async def run():
        async with httpx.AsyncClient(transport=httpx.MockTransport(handler)) as http_client:
            return await VisionClient(http_client=http_client).analyze_image(
                b"image-bytes",
                mime_type="image/png",
                context=VisionContext(
                    filename="chart.png",
                    file_type=DocumentFileType.IMAGE,
                    page_number=1,
                    source_element=SourceElementType.IMAGE_FILE,
                ),
            )

    result = __import__("asyncio").run(run())

    assert result.status is VisionEnrichmentStatus.SUCCESS
    assert result.supplemental_text == "Final chart label"


def test_document_parser_reports_file_page_count_not_block_count():
    import fitz

    document = fitz.open()
    document.new_page().insert_text((72, 72), "first page has text")
    document.new_page()
    document.new_page().insert_text((72, 72), "third page has text")
    content = document.tobytes()
    document.close()

    result = __import__("asyncio").run(
        DocumentParser().parse_document(content=content, filename="three-pages.pdf", mime_type="application/pdf")
    )

    assert result.page_count == 3
    assert result.block_count == 2
