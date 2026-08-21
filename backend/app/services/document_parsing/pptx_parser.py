from __future__ import annotations

import asyncio
import io
import os
from dataclasses import dataclass
from typing import Any

from .image_parser import ImageParser
from .models import (
    BlockStyleSignals,
    BoundingBox,
    DocumentBlock,
    DocumentBlockType,
    DocumentFileType,
    DocumentParsingProfile,
    ProcessingMode,
    SourceElementType,
)


@dataclass(frozen=True)
class _ShapeBlock:
    block: DocumentBlock
    top: float
    left: float
    height: float


class PPTXParser:
    def __init__(self, *, image_parser: ImageParser) -> None:
        self.image_parser = image_parser

    async def parse(
        self,
        *,
        content: bytes,
        filename: str,
        mime_type: str,
        profile: DocumentParsingProfile = DocumentParsingProfile.LEGACY_V2,
    ) -> list[DocumentBlock]:
        return await asyncio.to_thread(self._parse_sync, content, filename, mime_type, profile)

    async def page_count(self, *, content: bytes) -> int:
        return await asyncio.to_thread(_page_count, content)

    def _parse_sync(self, content: bytes, filename: str, mime_type: str, profile: DocumentParsingProfile) -> list[DocumentBlock]:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        try:
            presentation = Presentation(io.BytesIO(content))
        except Exception as exc:
            raise ValueError("pptx document could not be parsed") from exc
        if len(presentation.slides) > _int_env("DOCUMENT_MAX_PPT_SLIDES", 100):
            raise ValueError("pptx document exceeds slide limit")
        if profile is DocumentParsingProfile.LEGACY_V2:
            return self._parse_legacy(presentation, filename, mime_type, MSO_SHAPE_TYPE)
        return self._parse_structured(presentation, filename, mime_type, MSO_SHAPE_TYPE)

    def _parse_legacy(self, presentation: Any, filename: str, mime_type: str, shape_type: Any) -> list[DocumentBlock]:
        blocks: list[DocumentBlock] = []
        for page_number, slide in enumerate(presentation.slides, start=1):
            native_text = "\n".join(
                shape.text.strip()
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            )
            if native_text:
                blocks.append(DocumentBlock(
                    filename=filename, file_type=DocumentFileType.PPTX, page_number=page_number, block_index=1,
                    text=native_text, processing_mode=ProcessingMode.PPT_NATIVE_TEXT,
                    source_element=SourceElementType.PPT_TEXT_SHAPES,
                ))
            blocks.extend(self._image_blocks(
                presentation,
                slide,
                filename,
                page_number,
                shape_type,
                structured=False,
            ))
        return blocks

    def _parse_structured(self, presentation: Any, filename: str, mime_type: str, shape_type: Any) -> list[DocumentBlock]:
        del mime_type
        blocks: list[DocumentBlock] = []
        for page_number, slide in enumerate(presentation.slides, start=1):
            candidates: list[_ShapeBlock] = []
            title_shape = slide.shapes.title
            for source_index, shape in enumerate(slide.shapes, start=1):
                if getattr(shape, "has_table", False):
                    table_text = _table_markdown(shape.table)
                    if table_text:
                        candidates.append(_candidate(
                            DocumentBlock(
                                filename=filename, file_type=DocumentFileType.PPTX, page_number=page_number,
                                block_index=source_index, text=table_text, processing_mode=ProcessingMode.PPT_NATIVE_TEXT,
                                source_element=SourceElementType.PPT_TEXT_SHAPES, block_type=DocumentBlockType.TABLE,
                                bbox=_shape_bbox(shape), reading_order=source_index, structure_confidence=0.98,
                                source_element_index=source_index, table_header_rows=1,
                            ), shape,
                        ))
                    continue
                if not getattr(shape, "has_text_frame", False):
                    continue
                paragraphs = [paragraph for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()]
                for paragraph_index, paragraph in enumerate(paragraphs):
                    text = paragraph.text.strip()
                    is_title = (
                        paragraph_index == 0
                        and (shape is title_shape or (title_shape is None and float(shape.top) <= 0.8 * 914400))
                    )
                    block_type = DocumentBlockType.SLIDE_TITLE if is_title else (
                        DocumentBlockType.LIST_ITEM if paragraph.level > 0 else DocumentBlockType.SLIDE_BODY
                    )
                    local_start = shape.text.find(text)
                    candidates.append(_candidate(
                        DocumentBlock(
                            filename=filename, file_type=DocumentFileType.PPTX, page_number=page_number,
                            block_index=source_index, text=text, processing_mode=ProcessingMode.PPT_NATIVE_TEXT,
                            source_element=SourceElementType.PPT_TEXT_SHAPES, block_type=block_type,
                            bbox=_shape_bbox(shape), reading_order=source_index,
                            structure_confidence=0.95 if is_title else 0.82,
                            style_signals=BlockStyleSignals(list_level=paragraph.level),
                            source_char_start=max(0, local_start),
                            source_char_end=max(0, local_start) + len(text),
                            source_element_index=source_index,
                        ),
                        shape,
                    ))
            ordered = _row_band_sort(candidates)
            for reading_order, item in enumerate(ordered, start=1):
                item.block.reading_order = reading_order
                blocks.append(item.block)
            blocks.extend(self._image_blocks(
                presentation,
                slide,
                filename,
                page_number,
                shape_type,
                structured=True,
            ))
        return blocks

    def _image_blocks(
        self,
        presentation: Any,
        slide: Any,
        filename: str,
        page_number: int,
        shape_type: Any,
        *,
        structured: bool,
    ) -> list[DocumentBlock]:
        blocks: list[DocumentBlock] = []
        for image_index, shape in enumerate((s for s in slide.shapes if s.shape_type == shape_type.PICTURE), start=1):
            image_blocks = asyncio.run(self.image_parser.parse(
                content=shape.image.blob, filename=filename, mime_type=shape.image.content_type,
                page_number=page_number, file_type=DocumentFileType.PPTX,
                processing_mode=ProcessingMode.PPT_OCR, source_element=SourceElementType.PPT_EMBEDDED_IMAGE,
                source_element_index=image_index,
                image_coverage_ratio=(shape.width * shape.height) / (presentation.slide_width * presentation.slide_height),
                structured=structured,
            ))
            for block in image_blocks:
                if structured:
                    block.block_type = DocumentBlockType.IMAGE_DESCRIPTION
                block.bbox = _shape_bbox(shape)
                block.structure_confidence = block.ocr_confidence
            blocks.extend(image_blocks)
        return blocks


def _candidate(block: DocumentBlock, shape: Any) -> _ShapeBlock:
    bbox = block.bbox
    return _ShapeBlock(block=block, top=bbox.y0 if bbox else float(shape.top), left=bbox.x0 if bbox else float(shape.left), height=bbox.y1 - bbox.y0 if bbox else float(shape.height))


def _row_band_sort(items: list[_ShapeBlock]) -> list[_ShapeBlock]:
    bands: list[list[_ShapeBlock]] = []
    for item in sorted(items, key=lambda value: (value.top, value.left)):
        threshold = max(8.0, item.height * 0.75)
        if bands and abs(item.top - sum(value.top for value in bands[-1]) / len(bands[-1])) <= threshold:
            bands[-1].append(item)
        else:
            bands.append([item])
    return [item for band in bands for item in sorted(band, key=lambda value: value.left)]


def _shape_bbox(shape: Any) -> BoundingBox:
    scale = 914400.0
    return BoundingBox(
        x0=float(shape.left) / scale, y0=float(shape.top) / scale,
        x1=float(shape.left + shape.width) / scale, y1=float(shape.top + shape.height) / scale,
    )


def _table_markdown(table: Any) -> str:
    rows = [[(cell.text or "").strip() for cell in row.cells] for row in table.rows]
    rows = [row for row in rows if any(row)]
    if not rows or len(rows) < 2:
        return ""
    lines = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join("---" for _ in rows[0]) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in rows[1:])
    return "\n".join(lines)


def _int_env(name: str, default: int) -> int:
    try:
        value = int(os.getenv(name, str(default)))
    except ValueError:
        return default
    return value if value > 0 else default


def _page_count(content: bytes) -> int:
    from pptx import Presentation
    try:
        presentation = Presentation(io.BytesIO(content))
    except Exception as exc:
        raise ValueError("pptx document could not be parsed") from exc
    return len(presentation.slides)
