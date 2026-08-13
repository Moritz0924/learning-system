from __future__ import annotations

import math
import asyncio
import io
import time
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence

from backend.app.domain.rag.chunking.v3.config import HybridChunkPolicy, SemanticChunkPolicy, SizeGuardPolicy
from backend.app.domain.rag.chunking.v3.domain import SemanticSegment, StructuralRegion
from backend.app.domain.rag.chunking.v3.relations import AdjacentRelationChecker
from backend.app.domain.rag.chunking.v3.semantic import SemanticChunker, _semantic_unit
from backend.app.domain.rag.chunking.v3.size_guard import SizeGuard
from backend.app.domain.rag.chunking.v3.structure import StructureAwareChunker
from backend.app.domain.rag.chunking.v3.threshold import AdaptiveThresholdPolicy
from backend.app.services.document_parsing.models import DocumentParsingProfile
from backend.app.services.document_parsing.parser import DocumentParser
from backend.app.services.document_parsing.text_parser import StructuredTextParser
from backend.app.services.embeddings import DeterministicEmbeddingClient
from backend.app.services.token_counting import TiktokenTokenCounter

from .chunking_v3 import (
    ChunkingDocument,
    ChunkingQuery,
    EvidenceAnchor,
    RetrievedChunk,
    map_chunk_to_anchors,
    score_ranked_chunks,
)


class FixedThresholdPolicy:
    def __init__(self, threshold: float) -> None:
        self.threshold_value = threshold

    def threshold(self, scores: Sequence[float]) -> float:
        del scores
        return self.threshold_value

    @staticmethod
    def select(score: float, threshold: float) -> bool:
        return score > threshold


@dataclass(frozen=True)
class IndexedChunk:
    chunk_id: str
    document_id: str
    content: str
    metadata: dict
    token_count: int
    vector: tuple[float, ...]


@dataclass(frozen=True)
class VariantIndex:
    variant: str
    chunks: tuple[IndexedChunk, ...]
    provider_identity: str
    model: str
    dimensions: int
    timings: dict[str, float | int]
    diagnostics: dict[str, int]


class _TimedEncoder:
    def __init__(self, client) -> None:
        self.client = client
        self.calls = 0
        self.texts = 0
        self.seconds = 0.0

    def embed_batch(self, texts):
        started = time.perf_counter()
        values = self.client.embed_batch(list(texts))
        self.seconds += time.perf_counter() - started
        self.calls += 1
        self.texts += len(texts)
        return values


def build_variant_index(
    documents: Sequence[tuple[ChunkingDocument, str]],
    *,
    variant: str,
    policy: HybridChunkPolicy | None = None,
    fixed_threshold: float | None = None,
) -> VariantIndex:
    policy = policy or HybridChunkPolicy()
    encoder = DeterministicEmbeddingClient()
    semantic_encoder = _TimedEncoder(encoder)
    timings: dict[str, float | int] = {
        "parser_latency_seconds": 0.0,
        "semantic_segmentation_latency_seconds": 0.0,
        "semantic_embedding_latency_seconds": 0.0,
        "index_embedding_latency_seconds": 0.0,
        "total_ingestion_latency_seconds": 0.0,
        "logical_embedding_texts": 0,
        "physical_embedding_api_calls": 0,
    }
    diagnostics = _empty_diagnostics()
    total_started = time.perf_counter()
    counter = TiktokenTokenCounter(policy.tokenizer_id)
    raw_chunks: list[tuple[str, str, str, dict]] = []
    for document, text in documents:
        for position, (content, metadata) in enumerate(
            chunk_document(
                text,
                filename=document.filename,
                variant=variant,
                policy=policy,
                fixed_threshold=fixed_threshold,
                timings=timings,
                semantic_encoder=semantic_encoder,
                diagnostics=diagnostics,
            ),
            start=1,
        ):
            raw_chunks.append((f"{variant}-{document.document_id}-{position}", document.document_id, content, metadata))
    index_started = time.perf_counter()
    vectors = encoder.embed_batch([content for _, _, content, _ in raw_chunks])
    timings["index_embedding_latency_seconds"] = time.perf_counter() - index_started
    timings["logical_embedding_texts"] = len(raw_chunks)
    timings["physical_embedding_api_calls"] = 1 if raw_chunks else 0
    timings["semantic_embedding_latency_seconds"] = semantic_encoder.seconds
    timings["total_ingestion_latency_seconds"] = time.perf_counter() - total_started
    chunks = [IndexedChunk(
        chunk_id=chunk_id,
        document_id=document_id,
        content=content,
        metadata=metadata,
        token_count=counter.count(content),
        vector=tuple(vector),
    ) for (chunk_id, document_id, content, metadata), vector in zip(raw_chunks, vectors)]
    return VariantIndex(
        variant=variant,
        chunks=tuple(chunks),
        provider_identity=encoder.provider_identity,
        model=encoder.model,
        dimensions=encoder.dimensions,
        timings=timings,
        diagnostics=diagnostics,
    )


def chunk_document(
    text: str,
    *,
    filename: str,
    variant: str,
    policy: HybridChunkPolicy,
    fixed_threshold: float | None = None,
    timings: dict[str, float | int] | None = None,
    semantic_encoder=None,
    diagnostics: dict[str, int] | None = None,
) -> list[tuple[str, dict]]:
    parser_started = time.perf_counter()
    if variant == "A":
        chunks = _legacy_ingestion_chunks(text, filename=filename)
        if timings is not None:
            timings["parser_latency_seconds"] += time.perf_counter() - parser_started
        return chunks
    structured_blocks = _fixture_blocks(filename, text, profile=DocumentParsingProfile.STRUCTURED_V3)
    if timings is not None:
        timings["parser_latency_seconds"] += time.perf_counter() - parser_started

    blocks = structured_blocks
    structure = StructureAwareChunker()
    regions = structure.build_regions(blocks)
    if diagnostics is not None:
        diagnostics["semantic_regions"] += sum(region.region_type == "text" for region in regions)
    if variant == "P":
        structured_text = "\n\n".join(block.text for block in blocks if block.text.strip())
        return _deterministic_length_chunks(structured_text, policy)

    size_guard = SizeGuard(
        token_counter=TiktokenTokenCounter(policy.tokenizer_id),
        policy=policy.size,
    )
    for_semantic = {
        "B": SemanticChunkPolicy(
            local_window_weight=1.0, adjacent_weight=0.0, relation_penalty_weight=0.0,
        ),
        "C": SemanticChunkPolicy(
            local_window_weight=1.0, adjacent_weight=0.0, relation_penalty_weight=0.0,
        ),
        "D": policy.semantic,
        "E": policy.semantic,
    }[variant]
    semantic = SemanticChunker(
        encoder=semantic_encoder or DeterministicEmbeddingClient(),
        relation_checker=AdjacentRelationChecker(),
        threshold_policy=(
            FixedThresholdPolicy(fixed_threshold)
            if variant == "D" and fixed_threshold is not None
            else AdaptiveThresholdPolicy(
                min_samples=for_semantic.min_boundary_samples,
                mad_multiplier=for_semantic.mad_multiplier,
            )
        ),
        policy=for_semantic,
        batch_size=policy.semantic_batch_size,
    )
    output: list[tuple[str, dict]] = []
    for region in regions:
        if variant == "B":
            segments = [SemanticSegment(
                units=tuple(_semantic_unit(unit, index) for index, unit in enumerate(region.units)),
            )]
        else:
            semantic_started = time.perf_counter()
            segmentation = semantic.split_with_trace(region)
            segments = list(segmentation.segments)
            if timings is not None:
                timings["semantic_segmentation_latency_seconds"] += time.perf_counter() - semantic_started
            if diagnostics is not None:
                trace = segmentation.trace.boundaries
                diagnostics["candidate_boundaries"] += len(trace)
                diagnostics["selected_boundaries"] += sum(boundary.selected for boundary in trace)
                diagnostics["relation_adjusted_boundaries"] += sum(
                    boundary.continuation_score > 0 for boundary in trace
                )
                diagnostics["adaptive_threshold_regions"] += int(any(
                    boundary.adaptive_threshold is not None for boundary in trace
                ))
                diagnostics["fixed_threshold_regions"] += int(
                    variant == "D" and bool(trace) and fixed_threshold is not None
                )
        for candidate in size_guard.apply(segments):
            if diagnostics is not None:
                action = candidate.metadata.get("size_guard", {}).get("action")
                diagnostics["tiny_merges"] += int(action == "tiny_merge")
                diagnostics["hard_fallbacks"] += int(action == "hard_fallback")
            boundaries = [boundary.__dict__ for boundary in candidate.boundaries]
            source_units = {
                unit.unit_id: unit
                for unit in region.units
                if unit.unit_id in candidate.source_unit_ids
            }
            output.append((candidate.content, {
                "chunk_schema_version": "v3",
                "chunking_strategy": variant,
                "source_unit_ids": list(candidate.source_unit_ids),
                "source_spans": [
                    {
                        "page": unit.page_number,
                        "block_index": unit.block_index,
                        "char_start": unit.metadata.get("source_char_start"),
                        "char_end": unit.metadata.get("source_char_end"),
                    }
                    for unit in source_units.values()
                ],
                "semantic": {"boundaries": boundaries},
                "size_guard": dict(candidate.metadata.get("size_guard", {})),
            }))
    return output


def _empty_diagnostics() -> dict[str, int]:
    return {
        "semantic_regions": 0,
        "adaptive_threshold_regions": 0,
        "fixed_threshold_regions": 0,
        "candidate_boundaries": 0,
        "selected_boundaries": 0,
        "relation_adjusted_boundaries": 0,
        "tiny_merges": 0,
        "hard_fallbacks": 0,
    }


def _deterministic_length_chunks(text: str, policy: HybridChunkPolicy) -> list[tuple[str, dict]]:
    counter = TiktokenTokenCounter(policy.tokenizer_id)
    paragraphs = [part.strip() for part in text.split("\n\n") if part.strip()]
    result: list[tuple[str, dict]] = []
    current: list[str] = []
    for paragraph in paragraphs:
        if counter.count(paragraph) > policy.size.max_tokens:
            if current:
                result.append(("\n\n".join(current), {"chunk_schema_version": "v3", "chunking_strategy": "P"}))
                current = []
            result.extend(
                (fragment, {"chunk_schema_version": "v3", "chunking_strategy": "P"})
                for fragment in _token_safe_text_fragments(
                    paragraph,
                    counter=counter,
                    max_tokens=policy.size.max_tokens,
                )
            )
            continue
        trial = "\n\n".join([*current, paragraph])
        if current and counter.count(trial) > policy.size.max_tokens:
            content = "\n\n".join(current)
            result.append((content, {"chunk_schema_version": "v3", "chunking_strategy": "P"}))
            current = [paragraph]
        else:
            current.append(paragraph)
    if current:
        result.append(("\n\n".join(current), {"chunk_schema_version": "v3", "chunking_strategy": "P"}))
    return result or [(text, {"chunk_schema_version": "v3", "chunking_strategy": "P"})]


def _token_safe_text_fragments(
    text: str,
    *,
    counter: TiktokenTokenCounter,
    max_tokens: int,
) -> list[str]:
    fragments: list[str] = []
    start = 0
    while start < len(text):
        low, high = start + 1, len(text)
        safe_end = start
        while low <= high:
            midpoint = (low + high) // 2
            if counter.count(text[start:midpoint]) <= max_tokens:
                safe_end = midpoint
                low = midpoint + 1
            else:
                high = midpoint - 1
        if safe_end == start:
            raise ValueError("tokenizer cannot produce a token-safe Variant P fragment")
        fragments.append(text[start:safe_end])
        start = safe_end
    return fragments


def _fixture_blocks(text_filename: str, text: str, *, profile: DocumentParsingProfile):
    suffix = Path(text_filename).suffix.lower()
    if suffix in {".md", ".markdown", ".txt"}:
        mime_type = "text/markdown" if suffix in {".md", ".markdown"} else "text/plain"
        return StructuredTextParser().parse(
            content=text.encode("utf-8"),
            filename=text_filename,
            mime_type=mime_type,
        ).blocks
    content, mime_type = _materialize_fixture_bytes(suffix, text)
    result = asyncio.run(DocumentParser().parse_document(
        content=content,
        filename=text_filename,
        mime_type=mime_type,
        profile=profile,
    ))
    return result.blocks


def _legacy_ingestion_chunks(text: str, *, filename: str) -> list[tuple[str, dict]]:
    from backend.app.application.document_service import _parse_document_content

    suffix = Path(filename).suffix.lower()
    if suffix in {".md", ".markdown", ".txt"}:
        content = text.encode("utf-8")
        mime_type = "text/markdown" if suffix in {".md", ".markdown"} else "text/plain"
    else:
        content, mime_type = _materialize_fixture_bytes(suffix, text)
    parsed = _parse_document_content(content, filename=filename, mime_type=mime_type)
    return [(chunk["content"], chunk["metadata"]) for chunk in parsed.chunks]


def _materialize_fixture_bytes(suffix: str, text: str) -> tuple[bytes, str]:
    if suffix == ".pdf":
        import fitz

        document = fitz.open()
        lines = text.splitlines()
        midpoint = max(2, len(lines) // 2)
        for page_number, page_lines in enumerate((lines[:midpoint], lines[midpoint:]), start=1):
            page = document.new_page(width=612, height=792)
            page.insert_text((48, 56), page_lines[0].lstrip("# ")[:100], fontsize=20)
            body = "\n".join(line for line in page_lines[1:] if line.strip())
            page.insert_textbox(fitz.Rect(48, 90, 564, 720), body, fontsize=11)
        payload = document.tobytes()
        document.close()
        return payload, "application/pdf"
    if suffix == ".pptx":
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        lines = [line for line in text.splitlines() if line.strip()]
        midpoint = max(2, len(lines) // 2)
        for slide_lines in (lines[:midpoint], lines[midpoint:]):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
            title.text = slide_lines[0].lstrip("# ")[:100]
            body = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
            body.text = "\n".join(slide_lines[1:])
        buffer = io.BytesIO()
        presentation.save(buffer)
        return buffer.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    raise ValueError(f"unsupported fixture suffix: {suffix}")


def rank_chunks(index: VariantIndex, query: str, *, top_n: int = 20) -> list[IndexedChunk]:
    query_vector = DeterministicEmbeddingClient().embed(query)
    ranked = sorted(
        index.chunks,
        key=lambda chunk: _cosine(query_vector, chunk.vector),
        reverse=True,
    )
    return ranked[:top_n]


def evaluate_query(
    index: VariantIndex,
    query: ChunkingQuery,
    *,
    anchors: Sequence[EvidenceAnchor],
) -> dict[str, object]:
    anchors_by_id = {anchor.anchor_id: anchor for anchor in anchors}
    ranked = []
    for chunk in rank_chunks(index, query.query):
        covered = map_chunk_to_anchors(
            document_id=chunk.document_id,
            content=chunk.content,
            metadata=chunk.metadata,
            anchors=anchors,
        )
        ranked.append(RetrievedChunk(
            chunk_id=chunk.chunk_id,
            document_id=chunk.document_id,
            content=chunk.content,
            token_count=chunk.token_count,
            covered_anchor_ids=covered,
        ))
    return score_ranked_chunks(query=query, ranked=ranked, anchors_by_id=anchors_by_id)


def _cosine(left: Sequence[float], right: Sequence[float]) -> float:
    left_norm = math.sqrt(sum(value * value for value in left))
    right_norm = math.sqrt(sum(value * value for value in right))
    if not left_norm or not right_norm:
        return 0.0
    return sum(a * b for a, b in zip(left, right)) / left_norm / right_norm
